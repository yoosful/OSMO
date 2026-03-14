#!/usr/bin/env python3
"""Generate a manifest-driven Nut Pouring cookbook reproduction report PPTX."""

from __future__ import annotations

import argparse
import json
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    import yaml
except ImportError as error:  # pragma: no cover - surfaced at runtime
    raise SystemExit(
        "PyYAML is required to load report manifests. Install it with: pip install pyyaml"
    ) from error


BG_DARK = RGBColor(0x14, 0x1B, 0x2D)
BG_MID = RGBColor(0x1E, 0x27, 0x3B)
NVIDIA_GREEN = RGBColor(0x76, 0xB9, 0x00)
TEXT_WHITE = RGBColor(0xF5, 0xF7, 0xFA)
TEXT_GRAY = RGBColor(0xC7, 0xCF, 0xDB)
ACCENT_BLUE = RGBColor(0x4F, 0xB3, 0xFF)
SUCCESS_GREEN = RGBColor(0x29, 0xCC, 0x7A)
WARNING_YELLOW = RGBColor(0xFF, 0xD1, 0x3D)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BASE_DIR = Path(__file__).resolve().parent
DEFAULT_MANIFEST = BASE_DIR / "cookbook" / "nut_pouring" / "report_manifest.example.yaml"


def fmt_duration(seconds: float | int) -> str:
    sec = int(round(float(seconds)))
    minutes, rem = divmod(sec, 60)
    hours, minutes = divmod(minutes, 60)
    if hours > 0:
        return f"{hours}h {minutes}m {rem}s"
    return f"{minutes}m {rem}s"


def fmt_gib(bytes_value: int) -> str:
    return f"{bytes_value / (1024**3):.2f} GiB"


def resolve_path(path_value: str | None, manifest_dir: Path) -> str | None:
    if not path_value:
        return None
    path = Path(path_value).expanduser()
    if not path.is_absolute():
        path = manifest_dir / path
    return str(path.resolve())


@dataclass
class Sample:
    label: str
    artifact: str
    image: str | None = None
    video: str | None = None
    note: str | None = None


@dataclass
class Stage:
    number: str
    title: str
    workflow: str
    duration_seconds: float
    start_time: str
    end_time: str
    status: str
    input_dataset: str
    input_artifact: str
    output_dataset: str
    output_artifact: str
    details: list[str]
    sample_pairs: list[tuple[Sample, Sample]] = field(default_factory=list)
    summary_input_note: str | None = None
    summary_output_note: str | None = None
    appendix_samples: list[tuple[str, str, str]] = field(default_factory=list)


@dataclass
class ReportManifest:
    title: str
    subtitle: str
    report_title: str
    environment: str
    report_date: str
    result: str
    scope_points: list[str]
    execution_notes: list[str]
    dataset_chain: list[str]
    instances: list[str]
    stage_table: list[tuple[str, str, str, str]]
    stages: list[Stage]
    embedded_videos: list[dict[str, str]]
    pitfalls: list[str]
    cost_control: list[str]
    summary_points: list[str]
    artifact_pack_dir: str | None = None


def _load_sample(raw: dict[str, Any], manifest_dir: Path) -> Sample:
    return Sample(
        label=raw["label"],
        artifact=raw["artifact"],
        image=resolve_path(raw.get("image"), manifest_dir),
        video=resolve_path(raw.get("video"), manifest_dir),
        note=raw.get("note"),
    )


def load_manifest(path: Path) -> ReportManifest:
    raw = yaml.safe_load(path.read_text(encoding="utf-8"))
    manifest_dir = path.parent

    stages: list[Stage] = []
    for raw_stage in raw["stages"]:
        sample_pairs = []
        for pair in raw_stage.get("sample_pairs", []):
            sample_pairs.append((_load_sample(pair["input"], manifest_dir), _load_sample(pair["output"], manifest_dir)))

        appendix_samples = []
        for pair in sample_pairs:
            appendix_samples.append((pair[0].label, pair[0].artifact, pair[1].artifact))

        stages.append(
            Stage(
                number=str(raw_stage["number"]),
                title=raw_stage["title"],
                workflow=raw_stage["workflow"],
                duration_seconds=float(raw_stage["duration_seconds"]),
                start_time=raw_stage["start_time"],
                end_time=raw_stage["end_time"],
                status=raw_stage.get("status", "COMPLETED"),
                input_dataset=raw_stage["input_dataset"],
                input_artifact=raw_stage["input_artifact"],
                output_dataset=raw_stage["output_dataset"],
                output_artifact=raw_stage["output_artifact"],
                details=raw_stage["details"],
                sample_pairs=sample_pairs,
                summary_input_note=raw_stage.get("summary_input_note"),
                summary_output_note=raw_stage.get("summary_output_note"),
                appendix_samples=appendix_samples,
            )
        )

    return ReportManifest(
        title=raw["title"],
        subtitle=raw["subtitle"],
        report_title=raw["report_title"],
        environment=raw["environment"],
        report_date=raw["report_date"],
        result=raw["result"],
        scope_points=raw["scope_points"],
        execution_notes=raw["execution_notes"],
        dataset_chain=raw["dataset_chain"],
        instances=raw["instances"],
        stage_table=[tuple(item) for item in raw["stage_table"]],
        stages=stages,
        embedded_videos=raw.get("embedded_videos", []),
        pitfalls=raw["pitfalls"],
        cost_control=raw["cost_control"],
        summary_points=raw["summary_points"],
        artifact_pack_dir=resolve_path(raw.get("artifact_pack_dir"), manifest_dir),
    )


class ReportBuilder:
    def __init__(self, manifest: ReportManifest):
        self.manifest = manifest
        self.presentation = Presentation()
        self.presentation.slide_width = SLIDE_WIDTH
        self.presentation.slide_height = SLIDE_HEIGHT

    def set_slide_bg(self, slide, color=BG_DARK):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(
        self,
        slide,
        left,
        top,
        width,
        height,
        text,
        font_size=18,
        color=TEXT_WHITE,
        bold=False,
        alignment=PP_ALIGN.LEFT,
        font_name="Calibri",
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        paragraph = tf.paragraphs[0]
        paragraph.text = text
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.font.name = font_name
        paragraph.alignment = alignment
        return tf

    def add_paragraph(
        self,
        text_frame,
        text,
        font_size=14,
        color=TEXT_WHITE,
        bold=False,
        font_name="Calibri",
        alignment=PP_ALIGN.LEFT,
        space_before=Pt(2),
    ):
        paragraph = text_frame.add_paragraph()
        paragraph.text = text
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.font.name = font_name
        paragraph.alignment = alignment
        paragraph.space_before = space_before
        return paragraph

    def add_panel(self, slide, left, top, width, height, title):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_MID
        shape.line.color.rgb = RGBColor(0x2F, 0x3A, 0x55)
        text_frame = shape.text_frame
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = Pt(16)
        paragraph.font.bold = True
        paragraph.font.color.rgb = NVIDIA_GREEN
        paragraph.font.name = "Calibri"
        return shape

    def add_image_or_note(self, slide, image_path, left, top, width, height, note):
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            text_frame = self.add_text_box(
                slide,
                left,
                top,
                width,
                height,
                note,
                12,
                TEXT_GRAY,
                False,
                PP_ALIGN.LEFT,
                "Consolas",
            )
            text_frame.vertical_anchor = 1

    def add_movie_or_note(self, slide, movie_path, poster_path, left, top, width, height, note):
        if movie_path and os.path.exists(movie_path):
            poster = poster_path if poster_path and os.path.exists(poster_path) else None
            slide.shapes.add_movie(
                movie_path,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
                poster_frame_image=poster,
                mime_type="video/mp4",
            )
        else:
            self.add_image_or_note(slide, None, left, top, width, height, note)

    def add_footer_line(self, slide):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = NVIDIA_GREEN
        line.line.fill.background()

    def build(self):
        self.add_title_slide()
        self.add_scope_slide()
        self.add_status_slide()
        self.add_instances_slide()
        for stage in self.manifest.stages:
            self.add_stage_summary_slide(stage)
            self.add_stage_samples_slide(stage)
        self.add_embedded_videos_slide()
        self.add_appendix_slide()
        self.add_pitfalls_slide()
        self.add_cost_control_slide()
        self.add_summary_slide()
        return self.presentation

    def add_title_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.5, 12, 0.6, self.manifest.title, 16, NVIDIA_GREEN, True)
        self.add_text_box(slide, 0.8, 1.3, 12, 1.1, self.manifest.subtitle, 44, TEXT_WHITE, True)
        self.add_text_box(slide, 0.8, 2.35, 12, 0.7, self.manifest.report_title, 34, TEXT_WHITE, True)
        meta = self.add_text_box(slide, 0.8, 3.35, 12, 2.4, "Standalone report for full Step 1 to Step 6 execution", 20, TEXT_GRAY)
        self.add_paragraph(meta, f"Environment: {self.manifest.environment}", 17, TEXT_GRAY)
        self.add_paragraph(meta, f"Date: {self.manifest.report_date}", 17, TEXT_GRAY)
        self.add_paragraph(meta, f"Result: {self.manifest.result}", 19, SUCCESS_GREEN, True)
        self.add_footer_line(slide)

    def add_scope_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.3, 12, 0.6, "Report Scope", 32, TEXT_WHITE, True)
        scope = self.add_text_box(slide, 0.8, 1.0, 12, 5.9, "This document is a complete standalone execution report.", 20, NVIDIA_GREEN, True)
        self.add_paragraph(scope, "What is covered:", 16, TEXT_WHITE, True)
        for point in self.manifest.scope_points:
            self.add_paragraph(scope, f"- {point}", 15, TEXT_WHITE)
        self.add_paragraph(scope, "Execution note:", 16, TEXT_WHITE, True)
        for note in self.manifest.execution_notes:
            self.add_paragraph(scope, f"- {note}", 15, SUCCESS_GREEN if "COMPLETED" in note else WARNING_YELLOW)
        self.add_footer_line(slide)

    def add_status_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.3, 12, 0.6, "End-to-End Workflow Status", 32, TEXT_WHITE, True)
        text_frame = self.add_text_box(slide, 0.8, 1.0, 12, 6.1, "Step summary", 18, TEXT_GRAY)
        self.add_paragraph(text_frame, "Step  Workflow ID                       Duration    Status", 15, NVIDIA_GREEN, True, "Consolas")
        self.add_paragraph(text_frame, "--------------------------------------------------------------", 13, TEXT_GRAY, False, "Consolas")
        for step, workflow_id, duration, status in self.manifest.stage_table:
            self.add_paragraph(
                text_frame,
                f"{step:<5} {workflow_id:<33} {duration:<10} {status}",
                14,
                SUCCESS_GREEN if status == "COMPLETED" else WARNING_YELLOW,
                False,
                "Consolas",
            )
        self.add_paragraph(text_frame, "", 8, TEXT_GRAY)
        self.add_paragraph(text_frame, "Dataset chain:", 16, TEXT_WHITE, True)
        for line in self.manifest.dataset_chain:
            self.add_paragraph(text_frame, line, 14, TEXT_WHITE, False, "Consolas")
        self.add_footer_line(slide)

    def add_instances_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.3, 12, 0.6, "Instance Profile", 32, TEXT_WHITE, True)
        details = self.add_text_box(slide, 0.8, 1.0, 12, 5.9, "Infrastructure used for this reproduction", 20, NVIDIA_GREEN, True)
        for item in self.manifest.instances:
            self.add_paragraph(details, f"- {item}", 16, TEXT_WHITE)
        if self.manifest.artifact_pack_dir:
            self.add_paragraph(details, f"Artifact pack: {self.manifest.artifact_pack_dir}", 14, ACCENT_BLUE, False, "Consolas")
        self.add_footer_line(slide)

    def add_stage_summary_slide(self, stage: Stage):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.6, 0.2, 12.2, 0.6, f"Step {stage.number}: {stage.title}", 30, TEXT_WHITE, True)
        meta = self.add_text_box(
            slide,
            0.6,
            0.82,
            12.2,
            0.65,
            f"Workflow: {stage.workflow}   |   Duration: {fmt_duration(stage.duration_seconds)}   |   Status: {stage.status}",
            14,
            ACCENT_BLUE,
            True,
            font_name="Consolas",
        )
        self.add_paragraph(meta, f"Run window: {stage.start_time} to {stage.end_time}", 12, TEXT_GRAY, False, "Consolas", space_before=Pt(0))
        self.add_panel(slide, 0.6, 1.55, 6.1, 3.45, "Input Example")
        self.add_panel(slide, 6.65, 1.55, 6.1, 3.45, "Output Example")
        input_tf = self.add_text_box(slide, 0.85, 1.9, 5.6, 1.0, f"Dataset: {stage.input_dataset}", 13, TEXT_WHITE, True)
        self.add_paragraph(input_tf, f"Artifact: {stage.input_artifact}", 12, TEXT_GRAY, False, "Consolas")
        output_tf = self.add_text_box(slide, 6.9, 1.9, 5.6, 1.0, f"Dataset: {stage.output_dataset}", 13, TEXT_WHITE, True)
        self.add_paragraph(output_tf, f"Artifact: {stage.output_artifact}", 12, TEXT_GRAY, False, "Consolas")
        first_pair = stage.sample_pairs[0] if stage.sample_pairs else None
        self.add_image_or_note(
            slide,
            first_pair[0].image if first_pair else None,
            0.85,
            2.78,
            5.6,
            1.95,
            stage.summary_input_note or "No visual input artifact captured for this stage.",
        )
        self.add_image_or_note(
            slide,
            first_pair[1].image if first_pair else None,
            6.9,
            2.78,
            5.6,
            1.95,
            stage.summary_output_note or "No visual output artifact captured for this stage.",
        )
        detail_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.15), Inches(12.15), Inches(1.9)
        )
        detail_shape.fill.solid()
        detail_shape.fill.fore_color.rgb = BG_MID
        detail_shape.line.color.rgb = RGBColor(0x2F, 0x3A, 0x55)
        detail_tf = detail_shape.text_frame
        detail_tf.clear()
        paragraph = detail_tf.paragraphs[0]
        paragraph.text = "Run details"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = NVIDIA_GREEN
        paragraph.font.bold = True
        paragraph.font.name = "Calibri"
        for item in stage.details:
            self.add_paragraph(detail_tf, f"- {item}", 13, TEXT_WHITE)
        self.add_footer_line(slide)

    def add_stage_samples_slide(self, stage: Stage):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.6, 0.2, 12.2, 0.6, f"Step {stage.number}: Sample Evidence", 28, TEXT_WHITE, True)
        self.add_text_box(slide, 0.6, 0.82, 12.2, 0.5, f"{stage.workflow} | richer sample input/output evidence", 14, ACCENT_BLUE, True, font_name="Consolas")
        grid = [
            (0.6, 1.45, 4.0, 2.15),
            (4.66, 1.45, 4.0, 2.15),
            (8.72, 1.45, 4.0, 2.15),
        ]
        pairs = stage.sample_pairs[:3]
        for index, (left, top, width, height) in enumerate(grid):
            self.add_panel(slide, left, top, width, height, pairs[index][0].label if index < len(pairs) else f"Sample {index + 1}")
            if index < len(pairs):
                input_sample, output_sample = pairs[index]
                self.add_image_or_note(slide, input_sample.image, left + 0.12, top + 0.4, 1.75, 0.98, input_sample.note or input_sample.artifact)
                self.add_image_or_note(slide, output_sample.image, left + 2.05, top + 0.4, 1.75, 0.98, output_sample.note or output_sample.artifact)
                caption = self.add_text_box(slide, left + 0.12, top + 1.45, 3.7, 0.55, f"Input: {input_sample.artifact}", 10, TEXT_GRAY, False, font_name="Consolas")
                self.add_paragraph(caption, f"Output: {output_sample.artifact}", 10, TEXT_GRAY, False, "Consolas")
            else:
                self.add_text_box(slide, left + 0.12, top + 0.55, 3.7, 1.2, "No sample captured for this slot.", 12, TEXT_GRAY, False, font_name="Consolas")
        notes = self.add_text_box(slide, 0.8, 4.0, 12.0, 2.7, "Why these samples matter", 16, NVIDIA_GREEN, True)
        self.add_paragraph(notes, f"- Stage {stage.number} now shows multiple examples instead of a single token sample.", 13, TEXT_WHITE)
        self.add_paragraph(notes, "- Matching input/output demos make augmentation or conversion changes inspectable.", 13, TEXT_WHITE)
        if stage.sample_pairs:
            self.add_paragraph(notes, f"- Captured samples: {', '.join(sample.label for sample, _ in stage.sample_pairs[:3])}", 13, TEXT_WHITE, False, "Consolas")
        self.add_footer_line(slide)

    def add_embedded_videos_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.3, 12, 0.6, "Embedded Video Artifacts", 32, TEXT_WHITE, True)
        self.add_text_box(slide, 0.8, 0.95, 12, 0.5, "Double-click a clip in slideshow mode to play.", 14, TEXT_GRAY)
        panel_positions = [(0.6, 1.45, 4.15, 4.95), (4.85, 1.45, 4.15, 4.95), (9.1, 1.45, 3.65, 4.95)]
        for config, (left, top, width, height) in zip(self.manifest.embedded_videos[:3], panel_positions):
            self.add_panel(slide, left, top, width, height, config["title"])
            self.add_movie_or_note(
                slide,
                config.get("video"),
                config.get("poster"),
                left + 0.23,
                top + 0.5,
                min(width - 0.45, 3.7),
                min(height - 1.4, 2.1),
                config.get("missing_note", "Embedded video asset is missing."),
            )
        mapping = self.add_text_box(slide, 0.8, 5.1, 12.0, 1.6, "Clip mapping", 16, NVIDIA_GREEN, True)
        for config in self.manifest.embedded_videos[:3]:
            self.add_paragraph(mapping, f"- {config['title']}: {config['artifact']}", 13, TEXT_WHITE, False, "Consolas")
        self.add_footer_line(slide)

    def add_appendix_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.3, 12, 0.6, "Evidence Appendix", 32, TEXT_WHITE, True)
        appendix = self.add_text_box(slide, 0.8, 1.0, 12.0, 5.9, "Captured sample mapping", 18, NVIDIA_GREEN, True)
        for stage in self.manifest.stages:
            self.add_paragraph(appendix, f"Step {stage.number}: {stage.workflow}", 14, ACCENT_BLUE, True, "Consolas")
            for label, input_artifact, output_artifact in stage.appendix_samples[:3]:
                self.add_paragraph(appendix, f"- {label}: {input_artifact} -> {output_artifact}", 12, TEXT_WHITE, False, "Consolas")
        self.add_footer_line(slide)

    def add_pitfalls_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.3, 12, 0.6, "Common Pitfalls", 32, TEXT_WHITE, True)
        pit = self.add_text_box(slide, 0.8, 1.0, 12, 5.9, "Frequent setup issues to check first:", 20, NVIDIA_GREEN, True)
        for item in self.manifest.pitfalls:
            self.add_paragraph(pit, f"- {item}", 16, WARNING_YELLOW)
        self.add_footer_line(slide)

    def add_cost_control_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.3, 12, 0.6, "Resource and Cost Control", 32, TEXT_WHITE, True)
        text_frame = self.add_text_box(slide, 0.8, 1.0, 12, 5.9, "Idle scale-down and cleanup", 20, SUCCESS_GREEN, True)
        for item in self.manifest.cost_control:
            self.add_paragraph(text_frame, f"- {item}", 15, TEXT_WHITE if "Automation script" not in item else ACCENT_BLUE, False, "Consolas" if "Automation script" in item or "Cron:" in item else "Calibri")
        self.add_footer_line(slide)

    def add_summary_slide(self):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])
        self.set_slide_bg(slide)
        self.add_text_box(slide, 0.8, 0.3, 12, 0.6, "Final Summary", 32, TEXT_WHITE, True)
        summary = self.add_text_box(slide, 0.8, 1.0, 12, 5.9, "Full nut_pouring cookbook pipeline reproduced successfully.", 22, SUCCESS_GREEN, True)
        for point in self.manifest.summary_points:
            self.add_paragraph(summary, point, 17, SUCCESS_GREEN)
        self.add_footer_line(slide)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate Nut Pouring pipeline report PPTX")
    parser.add_argument("-m", "--manifest", default=str(DEFAULT_MANIFEST), help="YAML or JSON report manifest")
    parser.add_argument("-o", "--output", default="Nut_Pouring_Pipeline_Report.pptx", help="Output PPTX path")
    parser.add_argument("--dump-manifest-json", action="store_true", help="Print the loaded manifest as JSON and exit")
    return parser.parse_args()


def main():
    args = parse_args()
    manifest_path = Path(args.manifest).expanduser().resolve()
    manifest = load_manifest(manifest_path)
    if args.dump_manifest_json:
        print(json.dumps(manifest, default=lambda value: value.__dict__, indent=2))
        return
    builder = ReportBuilder(manifest)
    presentation = builder.build()
    output_path = Path(args.output).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    print(f"Report saved to: {output_path}")
    print(f"Slides: {len(presentation.slides)}")


if __name__ == "__main__":
    main()
